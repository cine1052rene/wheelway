"""WheelWay - 자동화 툴(Dify/n8n/Make.com) 연계 타당성 검토 PPT 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1B, 0x2A, 0x4A)
BLUE = RGBColor(0x2F, 0x6F, 0xED)
GRAY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
GREEN = RGBColor(0x1E, 0x8E, 0x3E)
RED = RGBColor(0xC5, 0x22, 0x1F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def add_bg(slide, color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    slide.shapes._spTree.remove(shp._element)
    slide.shapes._spTree.insert(2, shp._element)
    return shp


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=NAVY,
                 align=PP_ALIGN.LEFT, font="맑은 고딕", anchor=None, line_spacing=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return tb


def add_footer(slide, page_no):
    add_textbox(slide, 0.5, 7.08, 8, 0.35, "WheelWay · 자동화 툴 연계 타당성 검토", size=10, color=GRAY)
    add_textbox(slide, 12.0, 7.08, 0.8, 0.35, str(page_no), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def bar(slide, left, top, width, height, color=BLUE):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def rounded_box(slide, left, top, width, height, fill=LIGHT, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    return shp


# ---------- Slide 1: Title ----------
s = add_slide()
add_bg(s, NAVY)
bar(s, 0, 3.35, 13.333, 0.04, BLUE)
add_textbox(s, 1, 2.35, 11.3, 1.0, "자동화 툴(Dify · n8n · Make.com) 연계 타당성 검토",
            size=34, bold=True, color=WHITE)
add_textbox(s, 1, 3.5, 11.3, 0.6, "WheelWay 접근성 길찾기 프로젝트 — 라우팅 엔진 성능 보강 논의 후속 검토",
            size=16, color=RGBColor(0xC7, 0xD2, 0xE8))
add_textbox(s, 1, 6.5, 6, 0.5, "2026-07-23", size=13, color=RGBColor(0x9A, 0xA7, 0xC7))

# ---------- Slide 2: 배경 및 질문 ----------
s = add_slide()
add_bg(s)
bar(s, 0, 0, 0.15, 7.5, BLUE)
add_textbox(s, 0.7, 0.45, 11, 0.7, "검토 배경 및 질문", size=26, bold=True, color=NAVY)
bar(s, 0.7, 1.15, 1.1, 0.05, BLUE)

box = rounded_box(s, 0.7, 1.6, 11.9, 1.5, fill=LIGHT)
tf = box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "선행 논의: 네트워크 데이터가 700개 이상 역으로 확장되면 현재 라우팅 엔진(Dijkstra 방식)의 성능 병목이 커짐"
for r in p.runs:
    r.font.size = Pt(15); r.font.color.rgb = NAVY; r.font.bold = True; r.font.name = "맑은 고딕"
p2 = tf.add_paragraph()
p2.text = "→ 병목 원인: ① 매 반복 큐 전체 정렬 ② 간선 전체 스캔(인접 리스트 미사용) ③ 경로 배열 통째 복사"
for r in p2.runs:
    r.font.size = Pt(13); r.font.color.rgb = GRAY; r.font.name = "맑은 고딕"

add_textbox(s, 0.7, 3.4, 11.5, 0.5, "이번 질문", size=18, bold=True, color=BLUE)
add_textbox(s, 0.7, 3.95, 11.5, 1.2,
            "“디파이(Dify), n8n, Make.com 같은 자동화 툴과 연계하면 개발/구동이 더 빨라질까?”",
            size=20, bold=True, color=NAVY)
add_textbox(s, 0.7, 5.3, 11.5, 1.4,
            "→ 세 도구의 성격을 구분하고, 현재 병목(라우팅 엔진 계산 / 네트워크 데이터 확장)과 향후 운영 단계"
            "(데이터 갱신 · 알림 · AI 안내)로 나누어 적합성을 검토함",
            size=14, color=GRAY)
add_footer(s, 2)

# ---------- Slide 3: 세 도구 소개 ----------
s = add_slide()
add_bg(s)
bar(s, 0, 0, 0.15, 7.5, BLUE)
add_textbox(s, 0.7, 0.45, 11, 0.7, "검토 대상 도구 개요", size=26, bold=True, color=NAVY)
bar(s, 0.7, 1.15, 1.1, 0.05, BLUE)

tools = [
    ("Dify", "LLM 앱 / 챗봇 / RAG 빌더", "자연어 대화·문서질의 레이어 구축용"),
    ("n8n", "오픈소스 워크플로 자동화 (iPaaS)", "API 연결 · ETL · 스케줄 잡 · 웹훅, 셀프호스팅"),
    ("Make.com", "GUI 기반 iPaaS", "SaaS 연결 · 알림 자동화에 강함, 실행량 과금"),
]
x = 0.7
w = 3.85
for name, role, desc in tools:
    card = rounded_box(s, x, 1.7, w, 3.0, fill=WHITE, line_color=RGBColor(0xDD, 0xE2, 0xEC))
    add_textbox(s, x + 0.25, 1.95, w - 0.5, 0.5, name, size=22, bold=True, color=BLUE)
    bar(s, x + 0.25, 2.55, 0.6, 0.035, BLUE)
    add_textbox(s, x + 0.25, 2.8, w - 0.5, 0.9, role, size=14, bold=True, color=NAVY)
    add_textbox(s, x + 0.25, 3.65, w - 0.5, 1.0, desc, size=12.5, color=GRAY)
    x += w + 0.28

add_textbox(s, 0.7, 5.1, 11.7, 0.8,
            "공통점: 세 도구 모두 “서비스를 연결하고 자동화”하는 도구이며, 앱 내부 알고리즘 계산을 대신"
            "하거나 가속하는 도구가 아님", size=14, bold=True, color=RED)
add_footer(s, 3)

# ---------- Slide 4: 적합성 매트릭스 ----------
s = add_slide()
add_bg(s)
bar(s, 0, 0, 0.15, 7.5, BLUE)
add_textbox(s, 0.7, 0.45, 11, 0.7, "영역별 적합성 판단", size=26, bold=True, color=NAVY)
bar(s, 0.7, 1.15, 1.1, 0.05, BLUE)

rows = [
    ("라우팅 엔진 계산\n(경로탐색 알고리즘)", "부적합", RED,
     "그래프 최소경로 계산 기능이 없음. 외부 API로 빼면 매 요청마다 왕복 지연 발생 → 오히려 느려짐"),
    ("네트워크 데이터 확장\n(역·연결 정규화)", "부적합", RED,
     "일회성 스크립트 작업. Node 스크립트가 더 빠르고 디버깅도 쉬움 — 자동화 트리거의 장점이 살지 않음"),
    ("데이터 갱신 · 운영\n(엘리베이터 상태, 알림)", "적합", GREEN,
     "주기적 공공 API 폴링 → 캐시 갱신, 고장 시 관리자 알림(카톡/슬랙/메일) 자동화에 강점"),
    ("자연어 안내 챗봇\n(향후 확장)", "적합(향후)", GREEN,
     "“휠체어로 강남→서울역 어떻게 가?” 같은 자연어 질의 응대 레이어로 Dify 활용 가능"),
]
top = 1.55
row_h = 1.28
for label, verdict, color, desc in rows:
    rounded_box(s, 0.7, top, 3.0, row_h - 0.15, fill=LIGHT)
    add_textbox(s, 0.85, top + 0.08, 2.7, row_h - 0.3, label, size=13, bold=True, color=NAVY,
                anchor=MSO_ANCHOR.MIDDLE)
    chip = rounded_box(s, 3.85, top + 0.28, 1.7, 0.55, fill=color)
    tf = chip.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = verdict; p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "맑은 고딕"
    add_textbox(s, 5.75, top + 0.02, 6.7, row_h - 0.15, desc, size=12.5, color=GRAY,
                anchor=MSO_ANCHOR.MIDDLE)
    top += row_h
add_footer(s, 4)

# ---------- Slide 5: 장단점 비교표 ----------
s = add_slide()
add_bg(s)
bar(s, 0, 0, 0.15, 7.5, BLUE)
add_textbox(s, 0.7, 0.45, 11, 0.7, "도구별 장단점 비교", size=26, bold=True, color=NAVY)
bar(s, 0.7, 1.15, 1.1, 0.05, BLUE)

table_shape = s.shapes.add_table(4, 3, Inches(0.7), Inches(1.6), Inches(11.9), Inches(4.6))
table = table_shape.table
table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(5.05)
table.columns[2].width = Inches(5.05)

headers = ["도구", "장점", "단점"]
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = WHITE; r.font.name = "맑은 고딕"

data = [
    ("Dify", "AI 챗봇 / 자연어 안내 레이어를 빠르게 구축 가능",
     "라우팅·데이터 확장 병목과는 무관, 별도 인프라·운영비 필요"),
    ("n8n", "셀프호스팅 무료, 데이터 갱신 cron 파이프라인 구성 용이",
     "서버 직접 운영 부담, 핵심 성능 병목 해결엔 기여 없음"),
    ("Make.com", "GUI 기반으로 알림·SaaS 연동을 초고속으로 구성",
     "실행량 기준 과금, 벤더 종속, 워크플로 왕복 지연 존재"),
]
for r_i, (name, pro, con) in enumerate(data, start=1):
    cells = [table.cell(r_i, 0), table.cell(r_i, 1), table.cell(r_i, 2)]
    cells[0].text = name
    cells[1].text = pro
    cells[2].text = con
    for c_i, cell in enumerate(cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if r_i % 2 == 1 else LIGHT
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if c_i == 0 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(12.5); r.font.name = "맑은 고딕"
                r.font.color.rgb = NAVY if c_i == 0 else GRAY
                r.font.bold = (c_i == 0)
        cell.margin_left = Inches(0.15); cell.margin_right = Inches(0.15)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
add_footer(s, 5)

# ---------- Slide 6: 결론 및 추천 로드맵 ----------
s = add_slide()
add_bg(s)
bar(s, 0, 0, 0.15, 7.5, BLUE)
add_textbox(s, 0.7, 0.45, 11, 0.7, "결론 및 추천 로드맵", size=26, bold=True, color=NAVY)
bar(s, 0.7, 1.15, 1.1, 0.05, BLUE)

concl = rounded_box(s, 0.7, 1.55, 11.9, 1.35, fill=NAVY)
tf = concl.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.18)
p = tf.paragraphs[0]
p.text = "“개발 속도”를 위해 지금 끌어오는 것은 역효과 — 엔진 최적화 + 데이터 확장은 코드로 직접 처리"
for r in p.runs:
    r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "맑은 고딕"
p2 = tf.add_paragraph()
p2.text = "자동화 툴은 “운영·부가기능 도구”로 이후 단계에서 검토"
for r in p2.runs:
    r.font.size = Pt(13.5); r.font.color.rgb = RGBColor(0xC7, 0xD2, 0xE8); r.font.name = "맑은 고딕"

steps = [
    ("1", "라우팅 엔진 최적화", "인접 리스트 + 우선순위 큐 도입 (약 1~1.5h)", BLUE),
    ("2", "수도권 네트워크 확장", "역·연결 데이터 정규화 (2~5h, 범위에 따라 상이)", BLUE),
    ("3", "운영 자동화 도입 검토", "엘리베이터 실시간 상태 · 고장 알림에 n8n/Make 활용", GREEN),
]
top = 3.2
for num, title, desc, color in steps:
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(top), Inches(0.55), Inches(0.55))
    circle.fill.solid(); circle.fill.fore_color.rgb = color
    circle.line.fill.background(); circle.shadow.inherit = False
    tf = circle.text_frame
    p = tf.paragraphs[0]; p.text = num; p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "맑은 고딕"
    add_textbox(s, 1.5, top - 0.02, 4.3, 0.6, title, size=16, bold=True, color=NAVY)
    add_textbox(s, 1.5, top + 0.42, 10.3, 0.5, desc, size=12.5, color=GRAY)
    top += 1.05

add_footer(s, 6)

out_path = r"C:\project\Wheelway\docs\reports\WheelWay_Automation_Tools_Review.pptx"
prs.save(out_path)
print("SAVED:", out_path)
